# -*- coding: utf-8 -*-
"""
Genera la presentación de la Simulación de Plataforma SaaS.
Basado en el paper y el informe de hallazgos.
Salida: presentacion_simulacion_saas.pptx

Uso:
  python generar_presentacion.py
"""

import os
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Error: Se requiere python-pptx. Ejecute: pip install python-pptx")
    exit(1)

# Directorio base del proyecto (donde está el script)
BASE_DIR = Path(__file__).parent.resolve()


def _img_path(rel_path: str) -> Path:
    """Ruta absoluta a una imagen. Verifica que exista."""
    p = BASE_DIR / rel_path
    if not p.exists():
        print(f"Advertencia: No se encuentra {p}")
    return p


def add_title_slide(prs, title: str, subtitle: str = ""):
    """Slide de título."""
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(13.33)
    height = Inches(1.5)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(0x4a, 0x4a, 0x6a)
        p2.space_before = Pt(12)
    return slide


def add_content_slide(prs, title: str, bullets: list):
    """Slide con título y bullets."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    # Título
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    # Bullets
    tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.33), Inches(5.5))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = b
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(8)
    return slide


def add_image_slide(prs, title: str, img_path: str, caption: str = ""):
    """Slide con título e imagen centrada."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    # Título
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    # Imagen (centrada, ancho máx 12")
    path = _img_path(img_path)
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(0.8), Inches(1.0), width=Inches(11.7))
    if caption:
        tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.33), Inches(0.5))
        tf2 = tx2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = caption
        p2.font.size = Pt(12)
        p2.font.italic = True
        p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    return slide


def add_two_images_slide(prs, title: str, img1: str, img2: str, cap1: str = "", cap2: str = ""):
    """Slide con dos imágenes lado a lado."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    w, h = Inches(5.8), Inches(3.5)
    p1 = _img_path(img1)
    p2 = _img_path(img2)
    if p1.exists():
        slide.shapes.add_picture(str(p1), Inches(0.5), Inches(1.0), width=w, height=h)
    if p2.exists():
        slide.shapes.add_picture(str(p2), Inches(6.9), Inches(1.0), width=w, height=h)
    if cap1:
        tx1 = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(5.8), Inches(0.4))
        tx1.text_frame.paragraphs[0].text = cap1
        tx1.text_frame.paragraphs[0].font.size = Pt(11)
        tx1.text_frame.paragraphs[0].font.italic = True
    if cap2:
        tx2 = slide.shapes.add_textbox(Inches(6.9), Inches(4.6), Inches(5.8), Inches(0.4))
        tx2.text_frame.paragraphs[0].text = cap2
        tx2.text_frame.paragraphs[0].font.size = Pt(11)
        tx2.text_frame.paragraphs[0].font.italic = True
    return slide


def add_three_images_slide(prs, title: str, imgs: list, captions: list = None):
    """Slide con tres imágenes en fila."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    w, h = Inches(3.8), Inches(2.8)
    captions = captions or [""] * len(imgs)
    for i, (img, cap) in enumerate(zip(imgs, captions)):
        left = Inches(0.5 + i * 4.2)
        path = _img_path(img)
        if path.exists():
            slide.shapes.add_picture(str(path), left, Inches(1.0), width=w, height=h)
        if cap:
            t = slide.shapes.add_textbox(left, Inches(3.95), Inches(3.8), Inches(0.4))
            t.text_frame.paragraphs[0].text = cap
            t.text_frame.paragraphs[0].font.size = Pt(10)
            t.text_frame.paragraphs[0].font.italic = True
    return slide


def add_questions_slide(prs):
    """Slide final de preguntas."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    tx = slide.shapes.add_textbox(Inches(2), Inches(2.8), Inches(9.33), Inches(1.2))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = "¿Preguntas?"
    p.font.size = Pt(48)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)  # 16:9

    # 1. Portada
    add_title_slide(
        prs,
        "Simulación de Plataforma SaaS de Soporte Técnico Bajo Demanda",
        "Benchmark de viabilidad económica a 10 años",
    )

    # 2. Integrantes
    add_content_slide(
        prs,
        "Integrantes",
        [
            "Claudio Hernán Luciano Sonntag (120580-8)",
            "Ariel Gastón Pacay (121237-0)",
            "Alison Michella Reynoso (156899-1)",
            "",
            "Cátedra: Simulación — Ing. Erica M. Milin, Ing. David Carlos Mammana",
            "UTN FRBA — Ingeniería en Sistemas de Información",
        ],
    )

    # 3. Negocio y necesidad
    add_content_slide(
        prs,
        "Negocio y necesidad",
        [
            "Plataforma que conecta técnicos con clientes (modelo SaaS)",
            "Evaluar viabilidad económica bajo distintas configuraciones",
            "Simulación permite explorar escenarios sin costos ni riesgos en entorno real",
        ],
    )

    # 4. Variables de control
    add_content_slide(
        prs,
        "Variables de control",
        [
            "N — Frecuencia de releases: Semanales (7d), Mensuales (30d), Trimestrales (90d)",
            "M — Presupuesto marketing: 500, 1.500, 2.500 créditos/mes",
            "AB Testing — Proporción Suscripción vs Prepago: 0-100, 50-50, 100-0",
        ],
    )

    # 5. Variables de estado y cortes
    add_content_slide(
        prs,
        "Variables de estado y cortes de control",
        [
            "Estado: T (reloj), clientes por modalidad, técnicos, créditos prepago",
            "Último día de implementación, contrataciones pendientes",
            "Cortes: día de equilibrio, mejor trimestre, métricas agregadas",
        ],
    )

    # 6. Modelo
    add_content_slide(
        prs,
        "Modelo de simulación",
        [
            "Delta T = 1 día — avance día a día",
            "Eventos discretos: Contratación/Renuncia, Llegada clientes, Atender, Implementación,",
            "Cobro suscripciones, Pago sueldos, Renovación prepagos",
        ],
    )

    # 7. Objetivos
    add_content_slide(
        prs,
        "Objetivos",
        [
            "Punto de equilibrio (día en que beneficio > 0)",
            "Beneficio neto acumulado a 10 años",
            "Satisfacción de clientes (suscripción vs prepago)",
        ],
    )

    # 8. Ejecución
    add_content_slide(
        prs,
        "Ejecución del benchmark",
        [
            "27 configuraciones × 5.000 corridas × 3.650 días",
            "Total: 135.000 escenarios simulados",
            "Casos relevantes: 3 configs × 1.000 corridas × 5-6 años",
        ],
    )

    # 9. Resultados globales (heatmap legible en lugar de comparacion_todas)
    add_image_slide(
        prs,
        "Resultados globales — Beneficio por AB × Releases × Marketing",
        "benchmark_10_anos/comparacion_heatmap_beneficio.png",
        "Beneficio final (M créditos). Verde = ganancia, Rojo = pérdida.",
    )

    add_image_slide(
        prs,
        "Top 10 y peores 10 configuraciones",
        "benchmark_10_anos/comparacion_top_bottom.png",
        "Configuraciones ordenadas por beneficio final",
    )

    add_image_slide(
        prs,
        "Tasa de equilibrio por configuración",
        "benchmark_10_anos/equilibrio_por_config.png",
        "% de corridas que alcanzaron beneficio > 0",
    )

    # 10. Zoom casos relevantes
    add_three_images_slide(
        prs,
        "Zoom: casos relevantes — Serie beneficio acumulado (media ± 1σ)",
        [
            "benchmark_10_anos/casos_5_anos/equilibrio_antes/serie_beneficio_acumulado.png",
            "benchmark_10_anos/casos_5_anos/cost_effective/serie_beneficio_acumulado.png",
            "benchmark_10_anos/casos_5_anos/menos_efectivo/serie_beneficio_acumulado.png",
        ],
        [
            "Equilibrio más rápido (0-100/Trim/MKT2500)",
            "Cost-effective (50-50/Trim/MKT500)",
            "Menos efectivo (100-0/Sem/MKT500)",
        ],
    )

    # 11. resultado_neto vs serie_beneficio
    add_two_images_slide(
        prs,
        "Corrida individual vs media — Cost-effective",
        "benchmark_10_anos/casos_5_anos/cost_effective/ejemplo/resultado_neto.png",
        "benchmark_10_anos/casos_5_anos/cost_effective/serie_beneficio_acumulado.png",
        "resultado_neto (1 corrida)",
        "serie_beneficio_acumulado (media 1000 corridas)",
    )

    # 12. Caso extremo
    add_image_slide(
        prs,
        "Caso extremo — MKT 10.000 créditos/mes",
        "graficos_benchmark/caso_extremo_mkt10000/serie_beneficio_acumulado.png",
        "50-50 / Trimestrales — Equilibrio ~7 meses (semana 28-35). 1000 corridas, 5 años.",
    )

    # 13. Diagrama adicional (no en paper)
    add_image_slide(
        prs,
        "Tiempo al equilibrio por frecuencia de releases",
        "benchmark_10_anos/dia_equilibrio_por_release.png",
        "Día medio en que se alcanza beneficio > 0",
    )

    # 14. Conclusiones
    add_content_slide(
        prs,
        "Conclusiones",
        [
            "La frecuencia de releases es el factor más determinante — semanales destruyen valor en todas las configs",
            "Releases mensuales o trimestrales permiten viabilidad (74%-100% equilibrio)",
            "Config cost-effective: 50-50 / Trimestrales / MKT 500 — ~242k créditos, 90% equilibrio",
            "La estabilidad de la plataforma es más importante que el modelo de negocio",
        ],
    )

    add_image_slide(
        prs,
        "Mejor vs peor configuración",
        "benchmark_10_anos/mejor_vs_peor_config.png",
        "Diferencia: 6,7 millones de créditos",
    )

    # 15. Preguntas
    add_questions_slide(prs)

    out_path = BASE_DIR / "presentacion_simulacion_saas.pptx"
    prs.save(str(out_path))
    print(f"Presentación generada: {out_path}")
    print(f"Total de diapositivas: {len(prs.slides)}")


if __name__ == "__main__":
    main()
